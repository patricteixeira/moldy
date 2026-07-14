"""Leitura defensiva de PPTX editado para o Document Graph do M3."""

from __future__ import annotations

import hashlib
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from brand_runtime.native.ooxml import OoxmlError, validate_ooxml
from brand_runtime.native.pptx import (
    _BODY_TYPES,
    _TITLE_TYPES,
    _first_text_run,
    _run_color,
    _run_font_family,
    _theme_style,
)
from brand_runtime.roundtrip.models import (
    BoundsPt,
    DocumentDiagnostic,
    DocumentGraph,
    DocumentNode,
    DocumentSource,
)


class PptxParseError(OoxmlError):
    """O arquivo não pode entrar no round-trip com segurança ou semântica."""


@dataclass(frozen=True, slots=True)
class ShapeIdentity:
    """Identidade semântica recuperada dos sinais nativos do shape."""

    role: str
    slot_id: str | None
    revision_id: str | None
    source: str


def _description_fields(shape) -> dict[str, str]:
    nodes = shape._element.xpath(".//p:cNvPr")
    if not nodes:
        return {}
    fields: dict[str, str] = {}
    for field in nodes[0].get("descr", "").split(";"):
        key, separator, value = field.partition("=")
        if separator and key and value:
            fields[key] = value
    return fields


def identify_shape(shape) -> ShapeIdentity | None:
    """Recupera a identidade estável usada pelo parser e pelo fixer."""
    fields = _description_fields(shape)
    if shape.name.startswith("br:"):
        parts = shape.name.split(":", 2)
        if len(parts) == 3 and parts[1]:
            return ShapeIdentity(
                parts[1],
                parts[2] or fields.get("slot"),
                fields.get("brand-revision"),
                "name",
            )

    if fields.get("brand-role"):
        return ShapeIdentity(
            fields["brand-role"],
            fields.get("slot"),
            fields.get("brand-revision"),
            "description",
        )

    if shape.is_placeholder:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type in _TITLE_TYPES:
            return ShapeIdentity("heading", None, None, "placeholder")
        if placeholder_type in _BODY_TYPES:
            return ShapeIdentity("body", None, None, "placeholder")
    return None


def _points(value) -> float:
    return round(value.pt, 4)


def parse_pptx_document_graph(path: Path) -> DocumentGraph:
    """Valida e converte os objetos semânticos de um PPTX em Document Graph."""
    diagnostics = validate_ooxml(path)
    blocking = [item for item in diagnostics if item.blocking]
    if blocking:
        raise PptxParseError(blocking[0].message)

    presentation = Presentation(path)
    nodes: list[DocumentNode] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        theme_colors, theme_fonts = _theme_style(slide)
        for shape in slide.shapes:
            identity = identify_shape(shape)
            if identity is None:
                continue
            run = _first_text_run(shape)
            kind = "picture" if shape.shape_type == MSO_SHAPE_TYPE.PICTURE else "text"
            nodes.append(
                DocumentNode(
                    id=f"slide-{slide_index}/shape-{shape.shape_id}",
                    slide_index=slide_index,
                    shape_id=shape.shape_id,
                    kind=kind,
                    name=shape.name,
                    role=identity.role,
                    slot_id=identity.slot_id,
                    brand_revision_id=identity.revision_id,
                    semantic_source=identity.source,
                    text=(shape.text if getattr(shape, "has_text_frame", False) else None),
                    font_family=(
                        _run_font_family(run, identity.role, theme_fonts)
                        if run is not None
                        else None
                    ),
                    font_size_pt=(run.font.size.pt if run is not None and run.font.size else None),
                    color=(_run_color(run, theme_colors) if run is not None else None),
                    bounds_pt=BoundsPt(
                        x=_points(shape.left),
                        y=_points(shape.top),
                        width=_points(shape.width),
                        height=_points(shape.height),
                    ),
                )
            )

    if not nodes:
        raise PptxParseError("O PPTX não contém objetos semânticos reconhecíveis.")

    with path.open("rb") as source_file:
        digest = hashlib.file_digest(source_file, "sha256").hexdigest()
    return DocumentGraph(
        source=DocumentSource(
            filename=path.name,
            sha256=digest,
            size_bytes=path.stat().st_size,
            slide_count=len(presentation.slides),
        ),
        nodes=nodes,
        diagnostics=[
            DocumentDiagnostic(
                code=item.code,
                severity=item.severity,
                message=item.message,
                part=item.part,
            )
            for item in diagnostics
        ],
    )
