"""Plano explicável e aplicação conservadora de correções em cópia PPTX."""

from __future__ import annotations

import hashlib
import os
import tempfile
from pathlib import Path
from typing import Any, Literal

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pydantic import Field

from brand_runtime.ir.models import BrandIR, CamelModel
from brand_runtime.native.ooxml import OoxmlError, validate_ooxml
from brand_runtime.native.pptx import _first_text_run
from brand_runtime.roundtrip.lint import RoundtripFinding, RoundtripReport, lint_roundtrip
from brand_runtime.roundtrip.models import BoundsPt, DocumentGraph
from brand_runtime.roundtrip.pptx import identify_shape, parse_pptx_document_graph

FixProperty = Literal["fontFamily", "fontSizePt", "color", "boundsPt"]


class RoundtripFixError(OoxmlError):
    """O plano ou o arquivo não permitem uma correção segura e reproduzível."""


class FixOperation(CamelModel):
    """Uma alteração mínima, rastreável até findings do linter."""

    id: str
    slide_index: int = Field(ge=1)
    node_id: str
    role: str
    slot_id: str | None = None
    property: FixProperty
    expected: Any
    source_codes: list[str]


class FixPlan(CamelModel):
    """Plano determinístico que nunca contém alterações de conteúdo textual."""

    schema_version: Literal["0.1.0"] = "0.1.0"
    baseline_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    edited_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    operations: list[FixOperation]
    deferred_finding_codes: list[str]


class FixResult(CamelModel):
    """Evidência do novo artefato e do relint executado depois da correção."""

    schema_version: Literal["0.1.0"] = "0.1.0"
    source_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    fixed_sha256: str = Field(pattern=r"^[0-9a-f]{64}$")
    output_filename: str
    applied_operation_ids: list[str]
    report: RoundtripReport


_PROPERTY_BY_CODE: dict[str, FixProperty] = {
    "font-changed": "fontFamily",
    "brand-font": "fontFamily",
    "font-size-changed": "fontSizePt",
    "brand-font-size": "fontSizePt",
    "color-changed": "color",
    "brand-color": "color",
    "geometry-changed": "boundsPt",
}
_SEVERITY_PRIORITY = {"info": 0, "warning": 1, "error": 2, "locked": 3}


def _font_size_target(finding: RoundtripFinding) -> float:
    if finding.code != "brand-font-size":
        if not isinstance(finding.expected, (int, float)):
            raise RoundtripFixError("O finding de tamanho não possui um alvo numérico.")
        return float(finding.expected)
    if not isinstance(finding.expected, dict):
        raise RoundtripFixError("A faixa de tamanho da marca está ausente.")
    minimum = finding.expected.get("minimumPt", finding.expected.get("minimum_pt"))
    maximum = finding.expected.get("maximumPt", finding.expected.get("maximum_pt"))
    if not isinstance(minimum, (int, float)) or not isinstance(maximum, (int, float)):
        raise RoundtripFixError("A faixa de tamanho da marca é inválida.")
    actual = finding.actual
    if not isinstance(actual, (int, float)):
        return float(maximum)
    return float(minimum if actual < minimum else maximum)


def _target(finding: RoundtripFinding, property_name: FixProperty) -> Any:
    if property_name == "fontSizePt":
        return _font_size_target(finding)
    if property_name == "boundsPt":
        try:
            return BoundsPt.model_validate(finding.expected)
        except (TypeError, ValueError) as error:
            raise RoundtripFixError(
                "O finding de geometria não possui uma caixa válida."
            ) from error
    if property_name in {"fontFamily", "color"} and not isinstance(finding.expected, str):
        raise RoundtripFixError(f"O finding «{finding.code}» não possui um alvo textual.")
    if property_name == "color" and (
        len(finding.expected) != 7 or not finding.expected.startswith("#")
    ):
        raise RoundtripFixError("O alvo de cor precisa usar o formato #RRGGBB.")
    return finding.expected


def build_fix_plan(edited: DocumentGraph, report: RoundtripReport) -> FixPlan:
    """Converte somente findings corrigíveis em operações deduplicadas."""
    if report.edited_sha256 != edited.source.sha256:
        raise RoundtripFixError("O relatório não pertence ao Document Graph editado informado.")
    nodes = {node.id: node for node in edited.nodes}
    candidates: dict[tuple[str, FixProperty], dict[str, Any]] = {}
    deferred: list[str] = []

    for finding in report.findings:
        property_name = _PROPERTY_BY_CODE.get(finding.code)
        if not finding.fixable or property_name is None:
            deferred.append(finding.code)
            continue
        if finding.node_id is None or finding.node_id not in nodes:
            raise RoundtripFixError(f"O finding «{finding.code}» não referencia um nó conhecido.")
        node = nodes[finding.node_id]
        key = (node.id, property_name)
        candidate = candidates.get(key)
        priority = _SEVERITY_PRIORITY[finding.severity]
        if candidate is None:
            candidates[key] = {
                "node": node,
                "property": property_name,
                "expected": _target(finding, property_name),
                "priority": priority,
                "source_codes": [finding.code],
            }
            continue
        candidate["source_codes"].append(finding.code)
        if priority > candidate["priority"]:
            candidate["expected"] = _target(finding, property_name)
            candidate["priority"] = priority

    ordered = sorted(
        candidates.values(),
        key=lambda item: (
            item["node"].slide_index,
            item["node"].slot_id or item["node"].role,
            item["property"],
        ),
    )
    operations = [
        FixOperation(
            id=f"op-{index:03d}",
            slide_index=item["node"].slide_index,
            node_id=item["node"].id,
            role=item["node"].role,
            slot_id=item["node"].slot_id,
            property=item["property"],
            expected=item["expected"],
            source_codes=sorted(set(item["source_codes"])),
        )
        for index, item in enumerate(ordered, start=1)
    ]
    return FixPlan(
        baseline_sha256=report.baseline_sha256,
        edited_sha256=report.edited_sha256,
        operations=operations,
        deferred_finding_codes=deferred,
    )


def _sha256(path: Path) -> str:
    with path.open("rb") as source_file:
        return hashlib.file_digest(source_file, "sha256").hexdigest()


def _resolve_shape(presentation: Presentation, operation: FixOperation):
    slide = presentation.slides[operation.slide_index - 1]
    matches = []
    for shape in slide.shapes:
        identity = identify_shape(shape)
        if identity is None or identity.role != operation.role:
            continue
        if operation.slot_id is not None and identity.slot_id != operation.slot_id:
            continue
        matches.append(shape)
    if len(matches) != 1:
        label = operation.slot_id or operation.role
        raise RoundtripFixError(
            f"O alvo semântico «{label}» não foi reencontrado de forma inequívoca."
        )
    return matches[0]


def _apply_operation(shape, operation: FixOperation) -> None:
    if operation.property == "boundsPt":
        bounds = BoundsPt.model_validate(operation.expected)
        shape.left = Pt(bounds.x)
        shape.top = Pt(bounds.y)
        shape.width = Pt(bounds.width)
        shape.height = Pt(bounds.height)
        return
    run = _first_text_run(shape)
    if run is None:
        raise RoundtripFixError(f"O alvo «{operation.node_id}» não possui texto formatável.")
    if operation.property == "fontFamily":
        run.font.name = str(operation.expected)
    elif operation.property == "fontSizePt":
        run.font.size = Pt(float(operation.expected))
    elif operation.property == "color":
        run.font.color.rgb = RGBColor.from_string(str(operation.expected).removeprefix("#"))


def apply_pptx_fix_plan(
    source: Path,
    out: Path,
    plan: FixPlan,
    baseline: DocumentGraph,
    ir: BrandIR | None = None,
) -> FixResult:
    """Aplica o plano em uma cópia, valida OOXML e executa o relint."""
    source = source.resolve(strict=True)
    out = out.resolve()
    if source.suffix.lower() != ".pptx":
        raise RoundtripFixError("O fixer desta versão aceita somente arquivos PPTX.")
    if source == out:
        raise RoundtripFixError("O arquivo editado nunca pode ser sobrescrito.")
    if plan.baseline_sha256 != baseline.source.sha256:
        raise RoundtripFixError("O plano não pertence ao baseline informado.")
    source_sha = _sha256(source)
    if plan.edited_sha256 != source_sha:
        raise RoundtripFixError("Os bytes do arquivo mudaram depois da geração do plano.")
    blocking = [item for item in validate_ooxml(source) if item.blocking]
    if blocking:
        raise RoundtripFixError(blocking[0].message)

    presentation = Presentation(source)
    if any(operation.slide_index > len(presentation.slides) for operation in plan.operations):
        raise RoundtripFixError("O plano referencia um slide inexistente.")
    for operation in plan.operations:
        _apply_operation(_resolve_shape(presentation, operation), operation)

    out.parent.mkdir(parents=True, exist_ok=True)
    handle, temp_name = tempfile.mkstemp(prefix=f".{out.stem}-", suffix=".pptx", dir=out.parent)
    os.close(handle)
    temp_path = Path(temp_name)
    try:
        presentation.save(temp_path)
        blocking = [item for item in validate_ooxml(temp_path) if item.blocking]
        if blocking:
            raise RoundtripFixError(blocking[0].message)
        fixed = parse_pptx_document_graph(temp_path)
        report = lint_roundtrip(baseline, fixed, ir)
        os.replace(temp_path, out)
    finally:
        temp_path.unlink(missing_ok=True)

    return FixResult(
        source_sha256=source_sha,
        fixed_sha256=fixed.source.sha256,
        output_filename=out.name,
        applied_operation_ids=[operation.id for operation in plan.operations],
        report=report,
    )
