"""Renderer e inspeção PPTX template-first para o Gate 0 do M2."""

from __future__ import annotations

import os
import tempfile
from contextlib import contextmanager
from math import cos, pi, sin
from dataclasses import dataclass
from pathlib import Path

from lxml import etree
import fitz
from PIL import Image, ImageColor, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.slide import SlideLayout
from pptx.util import Emu, Pt

from brand_runtime.ir.models import BrandIR, SemanticRole
from brand_runtime.kit.models import (
    Area,
    AssetLayer,
    ContentSpec,
    EditorArea,
    ImageValue,
    LayerOverride,
    LayoutSpec,
    ShapeLayer,
    Slot,
    TextValue,
)
from brand_runtime.native.ooxml import OoxmlError, validate_ooxml

_TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
_BODY_TYPES = {
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.OBJECT,
    PP_PLACEHOLDER.SUBTITLE,
}
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


@dataclass(frozen=True, slots=True)
class SemanticShape:
    """Estado reencontrado no arquivo após edição externa."""

    role: str
    name: str
    kind: str
    text: str | None
    font_family: str | None
    font_size_pt: float | None
    color: str | None


class PptxRenderError(OoxmlError):
    """O template ou os contratos não permitem um PPTX nativo coerente."""


def _compatible_layout(layout: SlideLayout) -> bool:
    placeholder_types = {placeholder.placeholder_format.type for placeholder in layout.placeholders}
    return bool(placeholder_types & _TITLE_TYPES) and bool(placeholder_types & _BODY_TYPES)


def _select_layout(presentation: Presentation, requested_name: str | None) -> SlideLayout:
    if requested_name:
        for layout in presentation.slide_layouts:
            if layout.name == requested_name and _compatible_layout(layout):
                return layout
        raise PptxRenderError(
            f"O layout nativo «{requested_name}» não possui título e corpo editáveis."
        )
    for layout in presentation.slide_layouts:
        if _compatible_layout(layout):
            return layout
    raise PptxRenderError("O template não possui um layout com título e corpo editáveis.")


def _tag_shape(shape, *, role: str, revision_id: str, slot_id: str) -> None:
    name = f"br:{role}:{slot_id}"
    shape.name = name
    nodes = shape._element.xpath(".//p:cNvPr")
    if not nodes:
        raise PptxRenderError(f"Não foi possível etiquetar o shape do slot «{slot_id}».")
    nodes[0].set("name", name)
    nodes[0].set(
        "descr",
        f"brand-role={role};brand-revision={revision_id};slot={slot_id}",
    )


def _send_to_back(shape) -> None:
    """Mantém imagens de conteúdo atrás dos placeholders de texto e do logo."""
    shape_tree = shape._element.getparent()
    shape_tree.remove(shape._element)
    # Os dois primeiros filhos do spTree são propriedades do grupo, não shapes.
    shape_tree.insert(2, shape._element)


def _set_text(
    shape,
    value: TextValue,
    slot: Slot,
    role: SemanticRole,
    ir: BrandIR,
    override: LayerOverride | None = None,
) -> None:
    text = value.text
    text_transform = (
        override.text_transform if override and override.text_transform else slot.text_transform
    )
    if text_transform == "uppercase":
        text = text.upper()
    shape.text = text
    paragraph = shape.text_frame.paragraphs[0]
    if not paragraph.runs:
        run = paragraph.add_run()
        run.text = text
    else:
        run = paragraph.runs[0]
    font_token = override.font_token if override and override.font_token else role.font
    color_token = override.color_token if override and override.color_token else role.color
    font = ir.fonts[font_token]
    color = ir.colors[color_token]
    run.font.name = font.family
    font_size_px = (
        override.font_size_px
        if override and override.font_size_px
        else slot.font_size_px or role.max_size_px
    )
    run.font.size = Pt(font_size_px * 0.75)
    font_weight = (
        override.font_weight
        if override and override.font_weight
        else slot.font_weight or font.weight
    )
    font_style = (
        override.font_style if override and override.font_style else slot.font_style or font.style
    )
    run.font.bold = font_weight >= 600
    run.font.italic = font_style == "italic"
    run.font.color.rgb = RGBColor.from_string(color.value.removeprefix("#"))
    text_align = (
        override.text_align
        if override and override.text_align
        else slot.text_align
        if slot.text_align != "left"
        else None
    )
    if text_align is not None:
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }[text_align]
    line_height = override.line_height if override and override.line_height else slot.line_height
    if line_height is not None:
        paragraph.line_spacing = line_height
    letter_spacing = (
        override.letter_spacing_em
        if override and override.letter_spacing_em is not None
        else slot.letter_spacing_em
    )
    if letter_spacing:
        run_properties = run._r.get_or_add_rPr()
        spacing = round(letter_spacing * font_size_px * 0.75 * 100)
        run_properties.set("spc", str(spacing))


def _placeholder(slide, accepted: set[PP_PLACEHOLDER]):
    for placeholder in slide.placeholders:
        if placeholder.placeholder_format.type in accepted:
            return placeholder
    return None


def _resolve_asset(path: str, asset_root: Path | None) -> Path:
    candidate = Path(path)
    if not candidate.is_absolute() and asset_root is not None:
        candidate = asset_root / candidate
    candidate = candidate.resolve()
    if not candidate.is_file():
        raise PptxRenderError(f"O asset «{path}» não foi encontrado.")
    try:
        with Image.open(candidate) as image:
            image.verify()
    except (OSError, ValueError) as error:
        raise PptxRenderError(f"O asset «{path}» não é uma imagem raster válida.") from error
    return candidate


@contextmanager
def _picture_asset(path: str, asset_root: Path | None):
    """Entrega um raster aceito pelo OOXML, convertendo SVG local sem alterar a fonte."""
    candidate = Path(path)
    if not candidate.is_absolute() and asset_root is not None:
        candidate = asset_root / candidate
    candidate = candidate.resolve()
    if not candidate.is_file():
        raise PptxRenderError(f"O asset «{path}» não foi encontrado.")
    if candidate.suffix.lower() != ".svg":
        yield _resolve_asset(path, asset_root)
        return
    with tempfile.TemporaryDirectory(prefix="brandrt-svg-") as directory:
        output = Path(directory) / "asset.png"
        try:
            with fitz.open(candidate) as document:
                pixmap = document[0].get_pixmap(matrix=fitz.Matrix(2, 2), alpha=True)
                pixmap.save(output)
        except (RuntimeError, ValueError) as error:
            raise PptxRenderError(f"O SVG «{path}» não pôde ser rasterizado.") from error
        yield output


def _normalized_box(
    presentation: Presentation,
    layout: LayoutSpec,
    area: Area | EditorArea,
) -> tuple[Emu, ...]:
    x, y, width, height = area
    canvas_width = layout.canvas.width_px
    canvas_height = layout.canvas.height_px
    return (
        Emu(round(presentation.slide_width * x / canvas_width)),
        Emu(round(presentation.slide_height * y / canvas_height)),
        Emu(round(presentation.slide_width * width / canvas_width)),
        Emu(round(presentation.slide_height * height / canvas_height)),
    )


def _position_shape(
    shape,
    presentation: Presentation,
    layout: LayoutSpec,
    slot: Slot,
    override: LayerOverride | None = None,
) -> None:
    left, top, width, height = _normalized_box(
        presentation, layout, override.area if override and override.area else slot.area
    )
    shape.left = left
    shape.top = top
    shape.width = width
    shape.height = height
    if override is not None and override.rotation_deg is not None:
        shape.rotation = override.rotation_deg


def _apply_rotation(shape, override: LayerOverride | None) -> None:
    """Mantém no arquivo editável a rotação declarada pelo editor."""
    if override is not None and override.rotation_deg is not None:
        shape.rotation = override.rotation_deg


def _apply_opacity(shape, opacity: float) -> None:
    """Aplica transparência OOXML ao texto ou à imagem sem rasterizar o shape."""
    if opacity >= 1:
        return
    value = str(round(opacity * 100000))
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        for blip in shape._element.xpath(".//a:blip"):
            for previous in blip.xpath("./a:alphaModFix"):
                blip.remove(previous)
            etree.SubElement(blip, f"{{{_A_NS}}}alphaModFix", amt=value)
        return
    for solid_fill in shape._element.xpath(".//a:solidFill"):
        for previous in solid_fill.xpath("./a:alpha"):
            solid_fill.remove(previous)
        etree.SubElement(solid_fill, f"{{{_A_NS}}}alpha", val=value)


def _slot_override(content: ContentSpec, slot: Slot) -> LayerOverride | None:
    return content.overrides.get(slot.id)


def _relative_luminance(color: str) -> float:
    """Calcula luminância WCAG de uma cor canônica #RRGGBB."""
    channels = [int(color[index : index + 2], 16) / 255 for index in (1, 3, 5)]
    linear = [
        value / 12.92 if value <= 0.04045 else ((value + 0.055) / 1.055) ** 2.4
        for value in channels
    ]
    return linear[0] * 0.2126 + linear[1] * 0.7152 + linear[2] * 0.0722


def _automatic_logo_token(ir: BrandIR, background_token: str | None) -> str | None:
    """Escolhe a variante semântica que contrasta com o fundo efetivo."""
    if background_token is None or background_token not in ir.colors:
        return None
    candidate = (
        "logo.onLight"
        if _relative_luminance(ir.colors[background_token].value) >= 0.179
        else "logo.onDark"
    )
    return candidate if candidate in ir.assets else None


def _render_shape_layers(
    slide,
    presentation: Presentation,
    ir: BrandIR,
    layout: LayoutSpec,
    content: ContentSpec,
) -> None:
    """Materializa o subconjunto portátil de retângulos/círculos como shapes editáveis."""
    layers = sorted(
        (layer for layer in layout.locked_layers if isinstance(layer, ShapeLayer)),
        key=lambda layer: layer.z_index,
    )
    for layer in layers:
        override = content.overrides.get(layer.id)
        if override is not None and override.hidden:
            continue
        area = override.area if override and override.area else layer.area
        left, top, width, height = _normalized_box(presentation, layout, area)
        shape_type = MSO_SHAPE.OVAL if layer.shape == "circle" else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        _apply_rotation(shape, override)
        shape.line.fill.background()
        color_token = (
            override.color_token if override and override.color_token else layer.color_token
        )
        color = ir.colors[color_token]
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color.value.removeprefix("#"))
        _apply_opacity(
            shape,
            override.opacity if override and override.opacity is not None else layer.opacity,
        )
        _tag_shape(
            shape,
            role="layer",
            revision_id=ir.revision.id,
            slot_id=layer.id,
        )


def _render_asset_layers(
    slide,
    presentation: Presentation,
    ir: BrandIR,
    layout: LayoutSpec,
    content: ContentSpec,
    asset_root: Path | None,
    background_token: str | None,
) -> None:
    """Materializa assets estruturais como imagens nativas e editáveis no PPTX."""
    layers = sorted(
        (layer for layer in layout.locked_layers if isinstance(layer, AssetLayer)),
        key=lambda layer: layer.z_index,
    )
    for layer in layers:
        override = content.overrides.get(layer.id)
        if override is not None and override.hidden:
            continue
        asset_token = (
            content.asset_bindings.get(layer.id)
            or (
                _automatic_logo_token(ir, background_token)
                if layer.asset_token.startswith("logo.")
                else None
            )
            or layer.asset_token
        )
        area = override.area if override and override.area else layer.area
        left, top, width, height = _normalized_box(presentation, layout, area)
        with _picture_asset(ir.assets[asset_token].path, asset_root) as asset:
            picture = slide.shapes.add_picture(str(asset), left, top, width=width, height=height)
        _apply_rotation(picture, override)
        _apply_opacity(
            picture,
            override.opacity if override and override.opacity is not None else layer.opacity,
        )
        _tag_shape(
            picture,
            role="logo" if asset_token.startswith("logo.") else "asset",
            revision_id=ir.revision.id,
            slot_id=layer.id,
        )


def _text_slots(layout: LayoutSpec, content: ContentSpec) -> list[tuple[Slot, TextValue]]:
    values: list[tuple[Slot, TextValue]] = []
    for slot in layout.slots:
        value = content.values.get(slot.id)
        if slot.kind == "text" and isinstance(value, TextValue):
            values.append((slot, value))
    return values


def _validate_contracts(ir: BrandIR, layout: LayoutSpec, content: ContentSpec) -> None:
    if content.brand_revision_id != ir.revision.id:
        raise PptxRenderError("O Content Spec não pertence à revisão de marca informada.")
    if content.layout_id != layout.id:
        raise PptxRenderError("O Content Spec não pertence ao Layout Spec informado.")
    if (
        content.background_color_token is not None
        and content.background_color_token not in ir.colors
    ):
        raise PptxRenderError("O fundo da peça referencia uma cor ausente do Brand IR.")
    bindable_assets = {slot.id for slot in layout.slots if slot.kind == "logo"} | {
        layer.id for layer in layout.locked_layers if isinstance(layer, AssetLayer)
    }
    for element_id, asset_token in content.asset_bindings.items():
        if element_id not in bindable_assets:
            raise PptxRenderError("Asset bindings só podem referenciar logos ou assets conhecidos.")
        if asset_token not in ir.assets:
            raise PptxRenderError("O logo vinculado à peça não existe no Brand IR.")
    for slot, value in _text_slots(layout, content):
        if slot.role not in ir.roles:
            raise PptxRenderError(f"O papel semântico «{slot.role}» não existe no Brand IR.")
        role = ir.roles[slot.role]
        if role.font not in ir.fonts or role.color not in ir.colors:
            raise PptxRenderError(f"O papel semântico «{slot.role}» possui referência ausente.")


def _surface_png(ir: BrandIR, layout: LayoutSpec, content: ContentSpec, output: Path) -> Path:
    """Rasteriza somente a textura; texto, imagens e logo continuam nativos."""
    surface = content.surface
    if surface is None:
        raise PptxRenderError("A superfície procedural não foi informada.")
    token = ir.colors.get(surface.color_token)
    if token is None:
        raise PptxRenderError("A superfície referencia uma cor ausente da marca.")
    width, height = layout.canvas.width_px, layout.canvas.height_px
    rgba = (*ImageColor.getrgb(token.value), round(surface.opacity * 255))
    overlay = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    scale = max(4, round(surface.scale_px))
    weight = max(1, round(surface.weight_px))

    def line_field(angle: float, spacing: int = scale, line_weight: int = weight) -> Image.Image:
        diagonal = max(8, int((width * width + height * height) ** 0.5) * 2)
        field = Image.new("RGBA", (diagonal, diagonal), (0, 0, 0, 0))
        field_draw = ImageDraw.Draw(field)
        for y in range(-diagonal, diagonal * 2, max(4, spacing)):
            field_draw.line((0, y, diagonal, y), fill=rgba, width=max(1, line_weight))
        rotated = field.rotate(angle, resample=Image.Resampling.BICUBIC)
        left = max(0, (diagonal - width) // 2)
        top = max(0, (diagonal - height) // 2)
        return rotated.crop((left, top, left + width, top + height))

    def add_lines(*angles: float, spacing: int = scale, line_weight: int = weight) -> None:
        nonlocal overlay
        for angle in angles:
            overlay = Image.alpha_composite(
                overlay,
                line_field(angle, spacing=spacing, line_weight=line_weight),
            )

    if surface.kind == "technical-grid":
        for x in range(0, width + 1, scale):
            draw.line((x, 0, x, height), fill=rgba, width=weight)
        for y in range(0, height + 1, scale):
            draw.line((0, y, width, y), fill=rgba, width=weight)
    elif surface.kind == "point-field":
        radius = max(1, weight)
        for y in range(0, height + scale, scale):
            for x in range(0, width + scale, scale):
                draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=rgba)
    elif surface.kind == "concentric-rings":
        maximum = int((width * width + height * height) ** 0.5)
        center = (width // 2, height // 2)
        for radius in range(scale, maximum, scale):
            draw.ellipse(
                (
                    center[0] - radius,
                    center[1] - radius,
                    center[0] + radius,
                    center[1] + radius,
                ),
                outline=rgba,
                width=weight,
            )
    elif surface.kind == "paper-grain":
        # Distribuição determinística e irregular, sem RNG nem estado global.
        step = max(7, scale // 2)
        radius = max(1, weight)
        for row, y in enumerate(range(0, height + step, step)):
            for column, x in enumerate(range(0, width + step, step)):
                px = (x + (row * 17 + column * 7) % step) % width
                py = (y + (column * 13 + row * 5) % step) % height
                if (row * 3 + column * 5) % 4 == 0:
                    draw.ellipse((px - radius, py - radius, px + radius, py + radius), fill=rgba)
    elif surface.kind == "paper-fibers":
        add_lines(surface.angle_deg + 7, spacing=scale, line_weight=weight)
        add_lines(surface.angle_deg + 83, spacing=max(4, round(scale * 1.7)))
        add_lines(surface.angle_deg - 31, spacing=max(4, round(scale * 2.3)))
    elif surface.kind == "flecked-paper":
        radius = max(1, weight)
        step = max(8, scale)
        for row, y in enumerate(range(0, height + step, step)):
            for column, x in enumerate(range(0, width + step, step)):
                if (row * 7 + column * 11) % 5 > 1:
                    continue
                dx = (column * 19 + row * 3) % step
                dy = (row * 23 + column * 5) % step
                span = radius * (2 + (row + column) % 3)
                draw.ellipse(
                    (x + dx - span, y + dy - radius, x + dx + span, y + dy + radius), fill=rgba
                )
    elif surface.kind == "dry-brush":
        add_lines(surface.angle_deg, spacing=scale, line_weight=max(weight * 3, 2))
        add_lines(surface.angle_deg + 2, spacing=max(4, round(scale * 1.35)), line_weight=weight)
    elif surface.kind == "linear-rhythm":
        add_lines(surface.angle_deg)
    elif surface.kind == "scanlines":
        add_lines(0, spacing=max(4, round(scale * 0.35)))
    elif surface.kind == "diagonal-hatch":
        add_lines(surface.angle_deg or 45)
    elif surface.kind == "crosshatch":
        angle = surface.angle_deg or 45
        add_lines(angle, angle + 90)
    elif surface.kind == "woven":
        add_lines(0, 90, line_weight=max(1, round(weight * 1.6)))
        add_lines(0, 90, spacing=max(4, scale // 2), line_weight=max(1, weight // 2))
    elif surface.kind == "micro-grid":
        minor = max(4, scale // 4)
        for x in range(0, width + 1, minor):
            draw.line((x, 0, x, height), fill=rgba, width=weight)
        for y in range(0, height + 1, minor):
            draw.line((0, y, width, y), fill=rgba, width=weight)
        for x in range(0, width + 1, scale):
            draw.line((x, 0, x, height), fill=rgba, width=max(weight + 1, round(weight * 1.8)))
        for y in range(0, height + 1, scale):
            draw.line((0, y, width, y), fill=rgba, width=max(weight + 1, round(weight * 1.8)))
    elif surface.kind == "isometric-grid":
        add_lines(30, 90, 150)
    elif surface.kind == "halftone":
        radius_large = max(1, round(weight * 1.8))
        radius_small = max(1, weight)
        for row, y in enumerate(range(0, height + scale, scale)):
            offset = scale // 2 if row % 2 else 0
            radius = radius_small if row % 2 else radius_large
            for x in range(-scale, width + scale, scale):
                cx = x + offset
                draw.ellipse((cx - radius, y - radius, cx + radius, y + radius), fill=rgba)
    elif surface.kind == "checkerboard":
        cell = max(2, scale // 2)
        for row, y in enumerate(range(0, height + cell, cell)):
            for column, x in enumerate(range(0, width + cell, cell)):
                if (row + column) % 2 == 0:
                    draw.rectangle((x, y, x + cell, y + cell), fill=rgba)
    elif surface.kind == "topographic":
        centers = (
            (round(width * 0.18), round(height * 0.38)),
            (round(width * 0.82), round(height * 0.67)),
        )
        maximum = int((width * width + height * height) ** 0.5)
        for index, center in enumerate(centers):
            step = max(4, round(scale * (1 + index * 0.35)))
            for radius in range(step, maximum, step):
                vertical = round(radius * (0.62 + index * 0.08))
                draw.ellipse(
                    (
                        center[0] - radius,
                        center[1] - vertical,
                        center[0] + radius,
                        center[1] + vertical,
                    ),
                    outline=rgba,
                    width=weight,
                )
    elif surface.kind == "sunburst":
        center = (width // 2, height // 2)
        radius = int((width * width + height * height) ** 0.5)
        step_degrees = max(4, min(30, round(scale / 4)))
        for degrees in range(
            round(surface.angle_deg), 360 + round(surface.angle_deg), step_degrees
        ):
            radians = degrees * pi / 180
            draw.line(
                (
                    center[0],
                    center[1],
                    center[0] + cos(radians) * radius,
                    center[1] + sin(radians) * radius,
                ),
                fill=rgba,
                width=weight,
            )
    elif surface.kind == "waves":
        maximum = int((width * width + height * height) ** 0.5)
        for center_x in (0, width):
            for radius in range(scale, maximum, scale):
                draw.ellipse(
                    (
                        center_x - radius,
                        height // 2 - radius // 2,
                        center_x + radius,
                        height // 2 + radius // 2,
                    ),
                    outline=rgba,
                    width=weight,
                )
    elif surface.kind == "terrazzo":
        step = max(12, scale)
        for row, y in enumerate(range(0, height + step, step)):
            for column, x in enumerate(range(0, width + step, step)):
                if (row + column * 2) % 3 != 0:
                    continue
                dx = (row * 17 + column * 13) % step
                dy = (column * 19 + row * 11) % step
                span = max(2, weight * (2 + (row + column) % 3))
                draw.polygon(
                    (
                        (x + dx, y + dy - span),
                        (x + dx + span, y + dy + span),
                        (x + dx - span, y + dy + span // 2),
                    ),
                    fill=rgba,
                )
    else:
        raise PptxRenderError(f"Textura procedural desconhecida: {surface.kind}.")

    overlay.save(output, format="PNG", optimize=True)
    return output


def render_pptx(
    template_path: Path,
    output_path: Path,
    ir: BrandIR,
    layout: LayoutSpec,
    content: ContentSpec,
    *,
    asset_root: Path | None = None,
    native_layout_name: str | None = None,
) -> Path:
    """Preenche um slide nativo e preserva masters, layouts e theme do template."""
    template_path = template_path.resolve()
    output_path = output_path.resolve()
    if template_path == output_path:
        raise PptxRenderError("O template original nunca pode ser sobrescrito.")
    blocking = [item for item in validate_ooxml(template_path) if item.blocking]
    if blocking:
        raise PptxRenderError(f"O template possui {len(blocking)} erro(s) estrutural(is).")
    _validate_contracts(ir, layout, content)

    presentation = Presentation(template_path)
    native_layout = _select_layout(presentation, native_layout_name)
    slide = presentation.slides.add_slide(native_layout)
    background_color_token = content.background_color_token
    if background_color_token is None and layout.background.kind == "color":
        background_color_token = layout.background.color_token
    if background_color_token is not None:
        background = ir.colors.get(background_color_token)
        if background is not None:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor.from_string(background.value.removeprefix("#"))
    if content.surface is not None:
        with tempfile.TemporaryDirectory(prefix="brandrt-surface-") as directory:
            surface_path = _surface_png(ir, layout, content, Path(directory) / "surface.png")
            picture = slide.shapes.add_picture(
                str(surface_path),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
            picture.name = "br:surface"
            _send_to_back(picture)
    _render_shape_layers(slide, presentation, ir, layout, content)
    _render_asset_layers(
        slide,
        presentation,
        ir,
        layout,
        content,
        asset_root,
        background_color_token,
    )
    text_values = _text_slots(layout, content)
    if not text_values:
        raise PptxRenderError("O slide precisa de ao menos um slot de texto preenchido.")

    title_slot, title_value = text_values[0]
    title_shape = _placeholder(slide, _TITLE_TYPES)
    if title_shape is None:
        raise PptxRenderError("O layout selecionado perdeu o placeholder de título.")
    title_override = _slot_override(content, title_slot)
    _position_shape(title_shape, presentation, layout, title_slot, title_override)
    if title_override is None or not title_override.hidden:
        _set_text(
            title_shape, title_value, title_slot, ir.roles[title_slot.role], ir, title_override
        )
    else:
        title_shape.text = ""
    _apply_opacity(
        title_shape,
        title_override.opacity
        if title_override and title_override.opacity is not None
        else title_slot.opacity,
    )
    _tag_shape(
        title_shape,
        role=title_slot.role,
        revision_id=ir.revision.id,
        slot_id=title_slot.id,
    )

    body_shape = _placeholder(slide, _BODY_TYPES)
    if body_shape is None:
        raise PptxRenderError("O layout selecionado perdeu o placeholder de corpo.")
    if len(text_values) > 1:
        body_slot, body_value = text_values[1]
        body_override = _slot_override(content, body_slot)
        _position_shape(body_shape, presentation, layout, body_slot, body_override)
        if body_override is None or not body_override.hidden:
            _set_text(
                body_shape, body_value, body_slot, ir.roles[body_slot.role], ir, body_override
            )
        else:
            body_shape.text = ""
        _apply_opacity(
            body_shape,
            body_override.opacity
            if body_override and body_override.opacity is not None
            else body_slot.opacity,
        )
        _tag_shape(
            body_shape,
            role=body_slot.role,
            revision_id=ir.revision.id,
            slot_id=body_slot.id,
        )
    else:
        body_shape.text = ""
        _tag_shape(body_shape, role="body", revision_id=ir.revision.id, slot_id="body")

    for extra_slot, extra_value in text_values[2:]:
        extra_override = _slot_override(content, extra_slot)
        area = extra_override.area if extra_override and extra_override.area else extra_slot.area
        left, top, width, height = _normalized_box(presentation, layout, area)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        _apply_rotation(text_box, extra_override)
        if extra_override is None or not extra_override.hidden:
            _set_text(
                text_box, extra_value, extra_slot, ir.roles[extra_slot.role], ir, extra_override
            )
        _apply_opacity(
            text_box,
            extra_override.opacity
            if extra_override and extra_override.opacity is not None
            else extra_slot.opacity,
        )
        _tag_shape(
            text_box,
            role=extra_slot.role,
            revision_id=ir.revision.id,
            slot_id=extra_slot.id,
        )

    for image_slot in (slot for slot in layout.slots if slot.kind == "image"):
        value = content.values.get(image_slot.id)
        if not isinstance(value, ImageValue):
            continue
        image_override = _slot_override(content, image_slot)
        if image_override is not None and image_override.hidden:
            continue
        asset = _resolve_asset(value.path, asset_root)
        area = image_override.area if image_override and image_override.area else image_slot.area
        left, top, width, height = _normalized_box(presentation, layout, area)
        picture = slide.shapes.add_picture(str(asset), left, top, width=width, height=height)
        _apply_rotation(picture, image_override)
        _apply_opacity(
            picture,
            image_override.opacity
            if image_override and image_override.opacity is not None
            else image_slot.opacity,
        )
        _tag_shape(
            picture,
            role="image",
            revision_id=ir.revision.id,
            slot_id=image_slot.id,
        )
        _send_to_back(picture)

    logo_slot = next((slot for slot in layout.slots if slot.kind == "logo"), None)
    if logo_slot is not None:
        logo_override = _slot_override(content, logo_slot)
        if logo_override is not None and logo_override.hidden:
            logo_slot = None
    if logo_slot is not None:
        logo_override = _slot_override(content, logo_slot)
        composition_mode = (
            getattr(ir.composition_rules.modes, layout.composition_mode, None)
            if ir.composition_rules is not None and layout.composition_mode is not None
            else None
        )
        asset_token = (
            content.asset_bindings.get(logo_slot.id)
            or _automatic_logo_token(ir, background_color_token)
            or logo_slot.asset_token
            or (composition_mode.logo_asset_token if composition_mode is not None else None)
            or next(iter(ir.assets), None)
        )
        if asset_token is None or asset_token not in ir.assets:
            raise PptxRenderError(
                "O layout exige logo, mas o Brand IR não possui asset compatível."
            )
        area = logo_override.area if logo_override and logo_override.area else logo_slot.area
        left, top, width, height = _normalized_box(presentation, layout, area)
        with _picture_asset(ir.assets[asset_token].path, asset_root) as asset:
            picture = slide.shapes.add_picture(str(asset), left, top, width=width, height=height)
        _apply_rotation(picture, logo_override)
        _apply_opacity(
            picture,
            logo_override.opacity
            if logo_override and logo_override.opacity is not None
            else logo_slot.opacity,
        )
        _tag_shape(
            picture,
            role="logo",
            revision_id=ir.revision.id,
            slot_id=logo_slot.id,
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    handle, temp_name = tempfile.mkstemp(
        prefix=f".{output_path.stem}-", suffix=".pptx", dir=output_path.parent
    )
    os.close(handle)
    temp_path = Path(temp_name)
    try:
        presentation.save(temp_path)
        blocking = [item for item in validate_ooxml(temp_path) if item.blocking]
        if blocking:
            raise PptxRenderError(f"O PPTX gerado possui {len(blocking)} erro(s) estrutural(is).")
        os.replace(temp_path, output_path)
    finally:
        temp_path.unlink(missing_ok=True)
    return output_path


def _first_text_run(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                return run
    return None


def _theme_style(slide) -> tuple[dict[str, str], dict[str, str]]:
    """Resolve cores e fontes herdadas do tema usado pelo slide."""
    master = slide.slide_layout.slide_master
    theme_part = next(
        (
            relationship.target_part
            for relationship in master.part.rels.values()
            if relationship.reltype.endswith("/theme")
        ),
        None,
    )
    if theme_part is None:
        return {}, {}

    parser = etree.XMLParser(resolve_entities=False, no_network=True, load_dtd=False)
    theme_root = etree.fromstring(theme_part.blob, parser=parser)
    scheme = theme_root.find(
        ".//a:themeElements/a:clrScheme",
        {"a": _A_NS},
    )
    colors: dict[str, str] = {}
    if scheme is not None:
        for slot in scheme:
            color_nodes = list(slot)
            if not color_nodes:
                continue
            color_node = color_nodes[0]
            color_kind = etree.QName(color_node).localname
            value = (
                color_node.get("val")
                if color_kind == "srgbClr"
                else color_node.get("lastClr")
                if color_kind == "sysClr"
                else None
            )
            if value is not None and len(value) == 6:
                colors[etree.QName(slot).localname] = f"#{value.upper()}"

    master_root = etree.fromstring(master.part.blob, parser=parser)
    color_map = master_root.find(".//p:clrMap", {"p": _P_NS})
    if color_map is not None:
        for alias, target in color_map.attrib.items():
            resolved = colors.get(target)
            if resolved is not None:
                colors[alias] = resolved

    fonts: dict[str, str] = {}
    for role, xpath in {
        "major": ".//a:fontScheme/a:majorFont/a:latin",
        "minor": ".//a:fontScheme/a:minorFont/a:latin",
    }.items():
        font = theme_root.find(xpath, {"a": _A_NS})
        typeface = font.get("typeface") if font is not None else None
        if typeface:
            fonts[role] = typeface
    return colors, fonts


def _run_color(run, theme_colors: dict[str, str]) -> str | None:
    color_format = run.font.color
    if color_format.type == MSO_COLOR_TYPE.RGB:
        return f"#{color_format.rgb}"
    if color_format.type == MSO_COLOR_TYPE.SCHEME:
        color_node = getattr(color_format._color, "_xClr", None)
        scheme_name = color_node.get("val") if color_node is not None else None
        if scheme_name is None:
            return None
        return theme_colors.get(scheme_name, f"theme:{scheme_name}")
    return None


def _run_font_family(run, role: str, theme_fonts: dict[str, str]) -> str | None:
    family = run.font.name
    if family not in {None, "", "+mj-lt", "+mn-lt"}:
        return family
    theme_role = "major" if family == "+mj-lt" or role == "heading" else "minor"
    return theme_fonts.get(theme_role)


def _shape_role(shape) -> str | None:
    if shape.name.startswith("br:"):
        parts = shape.name.split(":", 2)
        if len(parts) == 3 and parts[1]:
            return parts[1]
    nodes = shape._element.xpath(".//p:cNvPr")
    if nodes:
        description = nodes[0].get("descr", "")
        for field in description.split(";"):
            key, separator, value = field.partition("=")
            if separator and key == "brand-role" and value:
                return value
    # LibreOffice deliberately rewrites the non-visual properties of native
    # placeholders. The placeholder kind is the third, editor-owned semantic
    # signal and survives that save cycle even when name/descr do not.
    if shape.is_placeholder:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type in _TITLE_TYPES:
            return "heading"
        if placeholder_type in _BODY_TYPES:
            return "body"
    return None


def inspect_semantic_shapes(path: Path) -> list[SemanticShape]:
    """Reencontra roles e formatação depois de um save externo ou automatizado."""
    if path.suffix.lower() != ".pptx" or not path.is_file():
        raise PptxRenderError("Informe um arquivo PPTX existente.")
    presentation = Presentation(path)
    findings: list[SemanticShape] = []
    for slide in presentation.slides:
        theme_colors, theme_fonts = _theme_style(slide)
        for shape in slide.shapes:
            role = _shape_role(shape)
            if role is None:
                continue
            run = _first_text_run(shape)
            color = _run_color(run, theme_colors) if run is not None else None
            findings.append(
                SemanticShape(
                    role=role,
                    name=shape.name,
                    kind=("picture" if shape.shape_type == MSO_SHAPE_TYPE.PICTURE else "text"),
                    text=(shape.text if getattr(shape, "has_text_frame", False) else None),
                    font_family=(
                        _run_font_family(run, role, theme_fonts) if run is not None else None
                    ),
                    font_size_pt=(run.font.size.pt if run is not None and run.font.size else None),
                    color=color,
                )
            )
    return findings
