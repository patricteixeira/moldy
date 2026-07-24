from __future__ import annotations

import hashlib
import json
import zipfile
from datetime import UTC, datetime
from pathlib import Path

import pytest
from docx import Document
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR_INDEX
from pptx.util import Inches, Pt
from typer.testing import CliRunner

from brand_runtime.cli import app
from brand_runtime.ir.models import (
    BrandIR,
    BrandInfo,
    ColorToken,
    Evidence,
    FontToken,
    LogoAsset,
    RevisionInfo,
    SemanticRole,
)
from brand_runtime.kit.models import (
    AssetLayer,
    Background,
    Canvas,
    ContentSpec,
    LayerOverride,
    LayoutSpec,
    ShapeLayer,
    Slot,
    SURFACE_KINDS,
    SurfaceStyle,
    TextValue,
)
from brand_runtime.native.docx import render_docx
from brand_runtime.native.ooxml import canonical_ooxml_manifest, validate_ooxml
from brand_runtime.native.pptx import _surface_png, inspect_semantic_shapes, render_pptx
from brand_runtime.roundtrip.fix import RoundtripFixError, apply_pptx_fix_plan, build_fix_plan
from brand_runtime.roundtrip.lint import lint_roundtrip
from brand_runtime.roundtrip.pptx import parse_pptx_document_graph
from brand_runtime.native.preview import render_native_preview
from brand_runtime.native.theme import derive_branded_template

FIXED = datetime(2026, 7, 14, 12, 0, tzinfo=UTC)
GOLDEN = Path(__file__).parent / "golden" / "native-ooxml.json"
RUNNER = CliRunner()


def _golden_parts(package_type: str, hashes: dict[str, str]) -> dict[str, str]:
    if package_type == "pptx":
        prefixes = (
            "ppt/presentation.xml",
            "ppt/_rels/presentation.xml.rels",
            "ppt/theme/",
            "ppt/slideMasters/",
            "ppt/slideLayouts/",
            "ppt/slides/",
            "ppt/media/",
        )
    else:
        prefixes = (
            "word/document.xml",
            "word/_rels/document.xml.rels",
            "word/styles.xml",
            "word/theme/",
            "word/media/",
        )
    return {name: digest for name, digest in hashes.items() if name.startswith(prefixes)}


@pytest.fixture
def native_brand(tmp_path: Path) -> BrandIR:
    logo_path = tmp_path / "logo.png"
    image = Image.new("RGBA", (180, 180), "#F6F0E5")
    draw = ImageDraw.Draw(image)
    draw.ellipse((18, 18, 162, 162), fill="#173F2C")
    draw.ellipse((60, 60, 120, 120), fill="#D4A72C")
    # Sem compressão, os mesmos pixels produzem bytes idênticos em Windows e Linux.
    image.save(logo_path, compress_level=0)
    digest = hashlib.sha256(logo_path.read_bytes()).hexdigest()
    evidence = [
        Evidence(
            source_type="manual-entry",
            path="tests/fixture",
            confidence=1,
            authoritative=True,
            confirmed_at=FIXED,
        )
    ]
    return BrandIR(
        schema_version="0.3.0",
        brand=BrandInfo(name="Marca Fixture"),
        revision=RevisionInfo(id="brandrev_native_fixture", created_at=FIXED),
        colors={
            "color.primary": ColorToken(value="#173F2C", evidence=evidence),
            "color.secondary": ColorToken(value="#D4A72C", evidence=evidence),
            "color.background": ColorToken(value="#F6F0E5", evidence=evidence),
            "color.text": ColorToken(value="#10231A", evidence=evidence),
        },
        fonts={
            "font.heading": FontToken(
                family="Georgia",
                weight=700,
                source="referenced-only",
                evidence=evidence,
            ),
            "font.body": FontToken(
                family="Arial",
                weight=400,
                source="referenced-only",
                evidence=evidence,
            ),
        },
        roles={
            "heading": SemanticRole(
                font="font.heading",
                color="color.text",
                min_size_px=36,
                max_size_px=64,
                line_height=1.05,
            ),
            "body": SemanticRole(
                font="font.body",
                color="color.text",
                min_size_px=16,
                max_size_px=24,
                line_height=1.4,
            ),
        },
        assets={
            "logo.primary": LogoAsset(
                path=str(logo_path),
                sha256=digest,
                format="png",
                evidence=evidence,
            )
        },
    )


@pytest.fixture
def slide_contracts(native_brand: BrandIR) -> tuple[LayoutSpec, ContentSpec]:
    layout = LayoutSpec(
        id="statement-post-1x1",
        profile="post-1x1",
        name_pt="Prova nativa",
        canvas=Canvas(width_px=1080, height_px=1080, safe_area_px=48),
        background=Background(kind="color", color_token="color.background"),
        slots=[
            Slot(
                id="headline",
                kind="text",
                role="heading",
                area=(64, 180, 800, 260),
            ),
            Slot(
                id="body",
                kind="text",
                role="body",
                area=(64, 480, 800, 300),
            ),
            Slot(
                id="logo",
                kind="logo",
                asset_token="logo.primary",
                area=(884, 884, 132, 132),
                fit="fixed",
            ),
        ],
    )
    content = ContentSpec(
        layout_id=layout.id,
        brand_revision_id=native_brand.revision.id,
        values={
            "headline": TextValue(text="A marca continua editável"),
            "body": TextValue(text="Texto, imagem, master e layout seguem nativos."),
        },
    )
    return layout, content


@pytest.fixture
def document_contracts(native_brand: BrandIR) -> tuple[LayoutSpec, ContentSpec]:
    layout = LayoutSpec(
        id="one-pager-doc-a4",
        profile="doc-a4",
        name_pt="Documento nativo",
        canvas=Canvas(width_px=794, height_px=1123, safe_area_px=76),
        background=Background(kind="color", color_token="color.background"),
        slots=[
            Slot(id="title", kind="text", role="heading", area=(76, 76, 642, 120)),
            Slot(id="body", kind="text", role="body", area=(76, 226, 642, 725)),
            Slot(
                id="logo",
                kind="logo",
                asset_token="logo.primary",
                area=(622, 951, 96, 96),
                fit="fixed",
            ),
        ],
    )
    content = ContentSpec(
        layout_id=layout.id,
        brand_revision_id=native_brand.revision.id,
        values={
            "title": TextValue(text="Documento de marca"),
            "body": TextValue(text="Um documento de fluxo contínuo com estilos semânticos."),
        },
    )
    return layout, content


@pytest.fixture
def pptx_template(tmp_path: Path) -> Path:
    path = tmp_path / "source-template.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    presentation.save(path)
    return path


@pytest.fixture
def docx_template(tmp_path: Path) -> Path:
    path = tmp_path / "source-template.docx"
    document = Document()
    document.add_paragraph("{{slot:title}}")
    document.add_paragraph("{{slot:body}}")
    document.add_paragraph("{{slot:logo}}")
    document.save(path)
    return path


def _render_pptx(
    tmp_path: Path,
    pptx_template: Path,
    native_brand: BrandIR,
    slide_contracts: tuple[LayoutSpec, ContentSpec],
    name: str = "native.pptx",
) -> Path:
    layout, content = slide_contracts
    branded_template = tmp_path / f"branded-{name}"
    derive_branded_template(pptx_template, branded_template, native_brand)
    return render_pptx(
        branded_template,
        tmp_path / name,
        native_brand,
        layout,
        content,
        native_layout_name="Title and Content",
    )


def test_pptx_template_fill_preserves_native_structure_and_source(
    tmp_path,
    pptx_template,
    native_brand,
    slide_contracts,
):
    original_hash = hashlib.sha256(pptx_template.read_bytes()).hexdigest()
    output = _render_pptx(
        tmp_path,
        pptx_template,
        native_brand,
        slide_contracts,
    )

    assert hashlib.sha256(pptx_template.read_bytes()).hexdigest() == original_hash
    assert not [item for item in validate_ooxml(output) if item.blocking]
    presentation = Presentation(output)
    assert len(presentation.slide_masters) >= 1
    assert len(presentation.slide_layouts) >= 2
    assert len(presentation.slides) == 1
    assert presentation.slides[0].slide_layout.name == "Title and Content"
    assert str(presentation.slides[0].background.fill.fore_color.rgb) == "F6F0E5"
    shapes = {shape.role: shape for shape in inspect_semantic_shapes(output)}
    assert shapes["heading"].text == "A marca continua editável"
    assert shapes["heading"].font_family == "Georgia"
    assert shapes["body"].kind == "text"
    assert shapes["logo"].kind == "picture"


def test_pptx_uses_instance_background_and_logo_binding(
    tmp_path,
    pptx_template,
    native_brand,
    slide_contracts,
):
    secondary_path = tmp_path / "logo-secondary.png"
    Image.new("RGBA", (180, 180), "#D4A72C").save(secondary_path, compress_level=0)
    secondary_digest = hashlib.sha256(secondary_path.read_bytes()).hexdigest()
    branded = native_brand.model_copy(deep=True)
    branded.assets["logo.secondary"] = LogoAsset(
        path=str(secondary_path),
        sha256=secondary_digest,
        format="png",
        evidence=branded.assets["logo.primary"].evidence,
    )
    layout, content = slide_contracts
    instance = content.model_copy(deep=True)
    instance.background_color_token = "color.primary"
    instance.asset_bindings = {"logo": "logo.secondary"}

    output = _render_pptx(
        tmp_path,
        pptx_template,
        branded,
        (layout, instance),
        "instance-bindings.pptx",
    )

    presentation = Presentation(output)
    slide = presentation.slides[0]
    assert str(slide.background.fill.fore_color.rgb) == "173F2C"
    logo = next(shape for shape in slide.shapes if shape.name == "br:logo:logo")
    assert hashlib.sha256(logo.image.blob).hexdigest() == secondary_digest


def test_pptx_exports_edited_structural_logo_as_native_picture(
    tmp_path,
    pptx_template,
    native_brand,
    slide_contracts,
):
    alternate_path = tmp_path / "logo-structural.png"
    Image.new("RGBA", (180, 180), "#D4A72C").save(alternate_path, compress_level=0)
    alternate_digest = hashlib.sha256(alternate_path.read_bytes()).hexdigest()
    branded = native_brand.model_copy(deep=True)
    branded.assets["logo.onLight"] = LogoAsset(
        path=str(alternate_path),
        sha256=alternate_digest,
        format="png",
        evidence=branded.assets["logo.primary"].evidence,
    )
    layout, content = slide_contracts
    authored_layout = layout.model_copy(deep=True)
    authored_layout.locked_layers.append(
        AssetLayer(
            id="brand-mark",
            asset_token="logo.primary",
            area=(84, 84, 120, 120),
            fit="contain",
            z_index=4,
        )
    )
    instance = content.model_copy(deep=True)
    instance.asset_bindings = {"brand-mark": "logo.onLight"}
    instance.overrides = {"brand-mark": LayerOverride(area=(120, 96, 180, 180), opacity=0.7)}

    output = _render_pptx(
        tmp_path,
        pptx_template,
        branded,
        (authored_layout, instance),
        "structural-logo.pptx",
    )

    presentation = Presentation(output)
    logo = next(
        shape for shape in presentation.slides[0].shapes if shape.name == "br:logo:brand-mark"
    )
    assert hashlib.sha256(logo.image.blob).hexdigest() == alternate_digest
    assert logo.left == round(presentation.slide_width * 120 / 1080)
    assert logo._element.xpath(".//a:alphaModFix")[0].get("amt") == "70000"


def test_pptx_selects_logo_variant_from_effective_background(
    tmp_path,
    pptx_template,
    native_brand,
    slide_contracts,
):
    branded = native_brand.model_copy(deep=True)
    for token, color in (("logo.onLight", "#151515"), ("logo.onDark", "#FAFAFA")):
        path = tmp_path / f"{token}.png"
        Image.new("RGBA", (180, 180), color).save(path, compress_level=0)
        branded.assets[token] = LogoAsset(
            path=str(path),
            sha256=hashlib.sha256(path.read_bytes()).hexdigest(),
            format="png",
            evidence=branded.assets["logo.primary"].evidence,
        )
    layout, content = slide_contracts
    instance = content.model_copy(deep=True)
    instance.background_color_token = "color.primary"

    output = _render_pptx(
        tmp_path,
        pptx_template,
        branded,
        (layout, instance),
        "automatic-logo-variant.pptx",
    )

    logo = next(
        shape for shape in Presentation(output).slides[0].shapes if shape.name == "br:logo:logo"
    )
    assert hashlib.sha256(logo.image.blob).hexdigest() == branded.assets["logo.onDark"].sha256


def test_pptx_applies_editor_geometry_typography_and_opacity_overrides(
    tmp_path,
    pptx_template,
    native_brand,
    slide_contracts,
):
    layout, content = slide_contracts
    edited_content = content.model_copy(deep=True)
    edited_content.overrides = {
        "headline": LayerOverride(
            area=(120, 220, 700, 280),
            rotation_deg=-12,
            opacity=0.55,
            font_token="font.body",
            font_size_px=80,
            font_weight=500,
            font_style="italic",
            color_token="color.secondary",
            line_height=1.2,
            letter_spacing_em=0.04,
            text_align="right",
            text_transform="uppercase",
        ),
        "logo": LayerOverride(
            area=(-180, 780, 1500, 900),
            rotation_deg=27,
            opacity=0.4,
        ),
    }
    edited_content.surface = SurfaceStyle(
        kind="technical-grid",
        color_token="color.secondary",
        opacity=0.12,
        scale_px=48,
        weight_px=1.5,
        angle_deg=0,
    )
    output = _render_pptx(
        tmp_path,
        pptx_template,
        native_brand,
        (layout, edited_content),
        "overrides.pptx",
    )

    presentation = Presentation(output)
    slide = presentation.slides[0]
    headline = next(shape for shape in slide.shapes if shape.name == "br:heading:headline")
    logo = next(shape for shape in slide.shapes if shape.name == "br:logo:logo")
    run = headline.text_frame.paragraphs[0].runs[0]

    assert headline.text == "A MARCA CONTINUA EDITÁVEL"
    assert headline.left == round(presentation.slide_width * 120 / 1080)
    assert headline.top == round(presentation.slide_height * 220 / 1080)
    assert headline.width == round(presentation.slide_width * 700 / 1080)
    # OOXML serializa a rotação negativa no intervalo equivalente de 0 a 360.
    assert headline.rotation == pytest.approx(348)
    assert run.font.name == "Arial"
    assert run.font.size.pt == pytest.approx(60)
    assert run.font.bold is False
    assert run.font.italic is True
    assert str(run.font.color.rgb) == "D4A72C"
    assert headline.text_frame.paragraphs[0].alignment == 3
    assert headline._element.xpath(".//a:alpha")[0].get("val") == "55000"

    assert logo.left == round(presentation.slide_width * -180 / 1080)
    assert logo.width == round(presentation.slide_width * 1500 / 1080)
    assert logo.rotation == pytest.approx(27)
    assert logo._element.xpath(".//a:alphaModFix")[0].get("amt") == "40000"
    surface = next(shape for shape in slide.shapes if shape.name == "br:surface")
    assert surface.shape_type == 13
    assert surface.width == presentation.slide_width
    assert not [item for item in validate_ooxml(output) if item.blocking]


def test_pptx_preserves_template_typography_and_shape_layers_as_editable_objects(
    tmp_path,
    pptx_template,
    native_brand,
    slide_contracts,
):
    layout, content = slide_contracts
    authored_layout = layout.model_copy(deep=True)
    headline = next(slot for slot in authored_layout.slots if slot.id == "headline")
    headline.font_size_px = 128
    headline.font_weight = 800
    headline.font_style = "italic"
    headline.line_height = 0.9
    headline.text_align = "center"
    authored_layout.locked_layers = [
        ShapeLayer(
            id="template-field",
            shape="rectangle",
            area=(40, 120, 960, 720),
            color_token="color.secondary",
            opacity=0.18,
            z_index=0,
        )
    ]

    output = _render_pptx(
        tmp_path,
        pptx_template,
        native_brand,
        (authored_layout, content),
        "template-style.pptx",
    )
    presentation = Presentation(output)
    slide = presentation.slides[0]
    title = next(shape for shape in slide.shapes if shape.name == "br:heading:headline")
    field = next(shape for shape in slide.shapes if shape.name == "br:layer:template-field")
    run = title.text_frame.paragraphs[0].runs[0]

    assert run.font.size.pt == pytest.approx(96)
    assert run.font.bold is True
    assert run.font.italic is True
    assert title.text_frame.paragraphs[0].alignment == 2
    assert str(field.fill.fore_color.rgb) == "D4A72C"
    assert field._element.xpath(".//a:alpha")[0].get("val") == "18000"


@pytest.mark.parametrize("kind", SURFACE_KINDS)
def test_pptx_surface_catalog_is_visible_and_deterministic(
    tmp_path,
    native_brand,
    slide_contracts,
    kind,
):
    layout, content = slide_contracts
    content.surface = SurfaceStyle(
        kind=kind,
        color_token="color.secondary",
        opacity=0.18,
        scale_px=42,
        weight_px=2,
        angle_deg=17,
    )
    first = _surface_png(native_brand, layout, content, tmp_path / f"{kind}-1.png")
    second = _surface_png(native_brand, layout, content, tmp_path / f"{kind}-2.png")

    with Image.open(first) as rendered:
        assert rendered.getbbox() is not None
    assert first.read_bytes() == second.read_bytes()


def test_pptx_round_trip_recovers_role_and_changed_formatting(
    tmp_path,
    pptx_template,
    native_brand,
    slide_contracts,
):
    output = _render_pptx(
        tmp_path,
        pptx_template,
        native_brand,
        slide_contracts,
    )
    edited = tmp_path / "edited.pptx"
    presentation = Presentation(output)
    title = next(
        shape for shape in presentation.slides[0].shapes if shape.name.startswith("br:heading:")
    )
    title.text = "Título alterado fora do renderer"
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.name = "Courier New"
    run.font.color.rgb = RGBColor(0xA1, 0x32, 0x20)
    # LibreOffice rewrites native placeholder names and drops their descr.
    # The inspector must still recover the semantic role from placeholder type.
    for shape in presentation.slides[0].shapes:
        if shape.is_placeholder:
            shape.name = f"PlaceHolder {shape.placeholder_format.idx + 1}"
            nodes = shape._element.xpath(".//p:cNvPr")
            if nodes:
                nodes[0].attrib.pop("descr", None)
    presentation.save(edited)

    shape = next(item for item in inspect_semantic_shapes(edited) if item.role == "heading")
    assert shape.text == "Título alterado fora do renderer"
    assert shape.font_family == "Courier New"
    assert shape.color == "#A13220"
    assert not [item for item in validate_ooxml(edited) if item.blocking]


def test_pptx_preserva_slot_de_imagem_como_picture_editavel(
    tmp_path,
    pptx_template,
    native_brand,
):
    logo = native_brand.assets["logo.primary"]
    layout = LayoutSpec(
        id="quote-post-1x1",
        profile="post-1x1",
        name_pt="Imagem nativa",
        canvas=Canvas(width_px=1080, height_px=1080, safe_area_px=48),
        background=Background(kind="image-slot"),
        slots=[
            Slot(id="quote", kind="text", role="heading", area=(64, 120, 720, 220)),
            Slot(id="photo", kind="image", area=(0, 420, 1080, 660), fit="fixed"),
        ],
    )
    content = ContentSpec(
        layout_id=layout.id,
        brand_revision_id=native_brand.revision.id,
        values={
            "quote": TextValue(text="Imagem substituível"),
            "photo": {"kind": "image", "path": logo.path, "sha256": logo.sha256},
        },
    )
    branded = tmp_path / "branded-image-template.pptx"
    derive_branded_template(pptx_template, branded, native_brand)
    output = render_pptx(
        branded,
        tmp_path / "native-image.pptx",
        native_brand,
        layout,
        content,
        native_layout_name="Title and Content",
    )

    shapes = inspect_semantic_shapes(output)
    assert any(shape.role == "image" and shape.kind == "picture" for shape in shapes)
    roles = [shape.role for shape in shapes]
    assert roles.index("image") < roles.index("heading")
    assert not [item for item in validate_ooxml(output) if item.blocking]


def test_pptx_round_trip_resolves_theme_color_after_powerpoint_save(
    tmp_path,
    pptx_template,
    native_brand,
    slide_contracts,
):
    output = _render_pptx(
        tmp_path,
        pptx_template,
        native_brand,
        slide_contracts,
    )
    edited = tmp_path / "edited-theme-color.pptx"
    presentation = Presentation(output)
    title = next(
        shape for shape in presentation.slides[0].shapes if shape.name.startswith("br:heading:")
    )
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.name = None
    run.font.color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_1
    presentation.save(edited)

    shape = next(item for item in inspect_semantic_shapes(edited) if item.role == "heading")
    assert shape.font_family == "Georgia"
    assert shape.color == "#173F2C"
    assert not [item for item in validate_ooxml(edited) if item.blocking]


def test_roundtrip_parser_builds_document_graph_after_google_style_save(
    tmp_path,
    pptx_template,
    native_brand,
    slide_contracts,
):
    output = _render_pptx(tmp_path, pptx_template, native_brand, slide_contracts)
    baseline = parse_pptx_document_graph(output)
    edited = tmp_path / "google-slides-edited.pptx"
    presentation = Presentation(output)
    for index, shape in enumerate(presentation.slides[0].shapes, start=1):
        if not shape.name.startswith("br:"):
            continue
        if shape.name.startswith("br:heading:"):
            shape.text = "Continua sim"
            run = shape.text_frame.paragraphs[0].runs[0]
            run.font.name = "Arial"
            run.font.size = Pt(48)
            run.font.color.rgb = RGBColor(0xE5, 0x79, 0x00)
        shape.name = f"Google Shape;{index};p13"
    presentation.save(edited)

    graph = parse_pptx_document_graph(edited)

    assert graph.schema_version == "0.1.0"
    assert graph.source.filename == edited.name
    assert graph.source.slide_count == 1
    assert len(graph.source.sha256) == 64
    assert [node.role for node in graph.nodes] == ["heading", "body", "logo"]
    heading = graph.nodes[0]
    assert heading.slot_id == "headline"
    assert heading.brand_revision_id == native_brand.revision.id
    assert heading.semantic_source == "description"
    assert heading.text == "Continua sim"
    assert heading.font_family == "Arial"
    assert heading.color == "#E57900"
    assert heading.bounds_pt.width > 0

    report = lint_roundtrip(baseline, graph, native_brand)
    codes = [finding.code for finding in report.findings]
    assert codes == [
        "text-changed",
        "font-changed",
        "color-changed",
        "brand-font",
        "brand-color",
    ]
    assert report.summary.status == "review"
    assert report.summary.info == 1
    assert report.summary.warning == 4
    assert report.summary.error == 0
    assert report.summary.fixable == 4

    source_bytes = edited.read_bytes()
    plan = build_fix_plan(graph, report)
    assert [operation.property for operation in plan.operations] == ["color", "fontFamily"]
    assert plan.operations[0].source_codes == ["brand-color", "color-changed"]
    assert plan.operations[1].source_codes == ["brand-font", "font-changed"]

    corrected = tmp_path / "corrected-copy.pptx"
    result = apply_pptx_fix_plan(edited, corrected, plan, baseline, native_brand)
    corrected_graph = parse_pptx_document_graph(corrected)

    assert edited.read_bytes() == source_bytes
    assert corrected_graph.source.sha256 != graph.source.sha256
    assert corrected_graph.nodes[0].text == "Continua sim"
    assert corrected_graph.nodes[0].font_family == "Georgia"
    expected_color = native_brand.colors[native_brand.roles["heading"].color].value
    assert corrected_graph.nodes[0].color == expected_color
    assert result.applied_operation_ids == ["op-001", "op-002"]
    assert [finding.code for finding in result.report.findings] == ["text-changed"]
    assert result.report.summary.fixable == 0
    assert result.report.summary.status == "review"
    assert not [item for item in validate_ooxml(corrected) if item.blocking]

    stale_plan = plan.model_copy(update={"edited_sha256": "0" * 64})
    with pytest.raises(RoundtripFixError, match="bytes do arquivo mudaram"):
        apply_pptx_fix_plan(
            edited,
            tmp_path / "must-not-exist.pptx",
            stale_plan,
            baseline,
            native_brand,
        )
    assert not (tmp_path / "must-not-exist.pptx").exists()

    baseline_path = tmp_path / "baseline.json"
    edited_graph_path = tmp_path / "edited.json"
    ir_path = tmp_path / "brand-ir.json"
    for path, model in (
        (baseline_path, baseline),
        (edited_graph_path, graph),
        (ir_path, native_brand),
    ):
        path.write_text(model.model_dump_json(by_alias=True), encoding="utf-8")
    report_path = tmp_path / "report.json"
    cli_result = RUNNER.invoke(
        app,
        [
            "roundtrip-lint",
            str(baseline_path),
            str(edited_graph_path),
            "--brand-ir",
            str(ir_path),
            "--out",
            str(report_path),
        ],
    )
    assert cli_result.exit_code == 0, cli_result.output
    plan_path = tmp_path / "fix-plan.json"
    cli_result = RUNNER.invoke(
        app,
        [
            "roundtrip-plan",
            str(edited_graph_path),
            str(report_path),
            "--out",
            str(plan_path),
        ],
    )
    assert cli_result.exit_code == 0, cli_result.output
    cli_corrected = tmp_path / "cli-corrected-copy.pptx"
    result_path = tmp_path / "fix-result.json"
    cli_result = RUNNER.invoke(
        app,
        [
            "roundtrip-fix",
            str(edited),
            str(baseline_path),
            str(plan_path),
            "--brand-ir",
            str(ir_path),
            "--out",
            str(cli_corrected),
            "--result-out",
            str(result_path),
        ],
    )
    assert cli_result.exit_code == 0, cli_result.output
    cli_evidence = json.loads(result_path.read_text(encoding="utf-8"))
    assert cli_evidence["report"]["summary"]["fixable"] == 0
    assert edited.read_bytes() == source_bytes


def test_docx_template_fill_uses_semantic_styles_and_native_image(
    tmp_path,
    docx_template,
    native_brand,
    document_contracts,
):
    layout, content = document_contracts
    branded_template = tmp_path / "branded-template.docx"
    derive_branded_template(docx_template, branded_template, native_brand)
    output = render_docx(
        branded_template,
        tmp_path / "native.docx",
        native_brand,
        layout,
        content,
    )

    assert not [item for item in validate_ooxml(output) if item.blocking]
    document = Document(output)
    title = next(
        paragraph for paragraph in document.paragraphs if paragraph.text == "Documento de marca"
    )
    body = next(
        paragraph for paragraph in document.paragraphs if paragraph.text.startswith("Um documento")
    )
    assert title.style.name == "Brand Heading"
    assert body.style.name == "Brand Body"
    assert len(document.inline_shapes) == 1
    assert not any("{{slot:" in paragraph.text for paragraph in document.paragraphs)


def test_canonical_ooxml_goldens_are_deterministic(
    tmp_path,
    pptx_template,
    docx_template,
    native_brand,
    slide_contracts,
    document_contracts,
):
    first_pptx = _render_pptx(
        tmp_path,
        pptx_template,
        native_brand,
        slide_contracts,
        "first.pptx",
    )
    second_pptx = _render_pptx(
        tmp_path,
        pptx_template,
        native_brand,
        slide_contracts,
        "second.pptx",
    )
    first_manifest = canonical_ooxml_manifest(first_pptx)
    second_manifest = canonical_ooxml_manifest(second_pptx)
    assert first_manifest.part_hashes == second_manifest.part_hashes

    layout, content = document_contracts
    branded_docx = tmp_path / "branded.docx"
    derive_branded_template(docx_template, branded_docx, native_brand)
    output_docx = render_docx(
        branded_docx,
        tmp_path / "native.docx",
        native_brand,
        layout,
        content,
    )
    actual = {
        "pptx": _golden_parts("pptx", first_manifest.part_hashes),
        "docx": _golden_parts("docx", canonical_ooxml_manifest(output_docx).part_hashes),
    }
    expected = json.loads(GOLDEN.read_text(encoding="utf-8"))
    assert actual == expected


def test_external_relationship_is_blocking(tmp_path, pptx_template):
    unsafe = tmp_path / "unsafe.pptx"
    with zipfile.ZipFile(pptx_template) as source, zipfile.ZipFile(unsafe, "w") as destination:
        for info in source.infolist():
            payload = source.read(info.filename)
            if info.filename == "_rels/.rels":
                payload = payload.replace(
                    b"</Relationships>",
                    b'<Relationship Id="rExternal" '
                    b'Type="http://example.com/external" '
                    b'Target="https://example.com" TargetMode="External"/>'
                    b"</Relationships>",
                )
            destination.writestr(info, payload)
    diagnostics = validate_ooxml(unsafe)
    assert any(item.code == "ooxml.external_relationship" and item.blocking for item in diagnostics)


def test_preview_failure_is_non_destructive_and_explicit(
    tmp_path,
    pptx_template,
    native_brand,
    slide_contracts,
):
    output = _render_pptx(
        tmp_path,
        pptx_template,
        native_brand,
        slide_contracts,
    )
    original_hash = hashlib.sha256(output.read_bytes()).hexdigest()
    result = render_native_preview(
        output,
        tmp_path / "preview",
        converter_path=tmp_path / "missing-soffice.exe",
    )
    assert not result.ok
    assert any(item.code == "preview.converter_unavailable" for item in result.diagnostics)
    assert hashlib.sha256(output.read_bytes()).hexdigest() == original_hash


def test_native_cli_exposes_theme_pptx_and_docx_product_slice(
    tmp_path,
    pptx_template,
    docx_template,
    native_brand,
    slide_contracts,
    document_contracts,
):
    ir_path = tmp_path / "brand-ir.json"
    ir_path.write_text(native_brand.model_dump_json(by_alias=True), encoding="utf-8")
    slide_layout, slide_content = slide_contracts
    doc_layout, doc_content = document_contracts
    paths = {}
    for name, model in {
        "slide-layout": slide_layout,
        "slide-content": slide_content,
        "doc-layout": doc_layout,
        "doc-content": doc_content,
    }.items():
        path = tmp_path / f"{name}.json"
        path.write_text(model.model_dump_json(by_alias=True), encoding="utf-8")
        paths[name] = path

    themed_pptx = tmp_path / "themed.pptx"
    result = RUNNER.invoke(
        app,
        ["native-theme", str(ir_path), str(pptx_template), "--out", str(themed_pptx)],
    )
    assert result.exit_code == 0, result.output
    result = RUNNER.invoke(
        app,
        [
            "native-pptx",
            str(ir_path),
            str(paths["slide-layout"]),
            str(paths["slide-content"]),
            str(themed_pptx),
            "--assets-dir",
            str(tmp_path),
            "--native-layout",
            "Title and Content",
            "--out",
            str(tmp_path / "cli.pptx"),
        ],
    )
    assert result.exit_code == 0, result.output
    result = RUNNER.invoke(app, ["native-inspect", str(tmp_path / "cli.pptx")])
    assert result.exit_code == 0, result.output
    inspection = json.loads(result.output)
    assert [shape["role"] for shape in inspection["semanticShapes"]] == [
        "heading",
        "body",
        "logo",
    ]
    graph_path = tmp_path / "document-graph.json"
    result = RUNNER.invoke(
        app,
        [
            "roundtrip-parse",
            str(tmp_path / "cli.pptx"),
            "--out",
            str(graph_path),
        ],
    )
    assert result.exit_code == 0, result.output
    graph = json.loads(graph_path.read_text(encoding="utf-8"))
    assert graph["schemaVersion"] == "0.1.0"
    assert [node["role"] for node in graph["nodes"]] == ["heading", "body", "logo"]

    themed_docx = tmp_path / "themed.docx"
    result = RUNNER.invoke(
        app,
        ["native-theme", str(ir_path), str(docx_template), "--out", str(themed_docx)],
    )
    assert result.exit_code == 0, result.output
    result = RUNNER.invoke(
        app,
        [
            "native-docx",
            str(ir_path),
            str(paths["doc-layout"]),
            str(paths["doc-content"]),
            str(themed_docx),
            "--assets-dir",
            str(tmp_path),
            "--out",
            str(tmp_path / "cli.docx"),
        ],
    )
    assert result.exit_code == 0, result.output
