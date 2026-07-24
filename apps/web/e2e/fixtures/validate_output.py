"""Valida dimensões e estrutura dos arquivos exportados no E2E."""

from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def main() -> None:
    """Valida os quatro formatos oferecidos pelo editor."""
    kind, raw_path = sys.argv[1:3]
    path = Path(raw_path)
    if kind in {"png", "png-4x5"}:
        with Image.open(path) as image:
            image.verify()
        with Image.open(path) as image:
            assert image.format == "PNG"
            expected_size = (1080, 1350) if kind == "png-4x5" else (1080, 1080)
            assert image.size == expected_size
        return

    if kind == "pdf":
        reader = PdfReader(path)
        assert len(reader.pages) == 1
        box = reader.pages[0].mediabox
        width = float(box.width)
        height = float(box.height)
        assert abs(width - 595) <= 2
        assert abs(height - 842) <= 2
        return

    if kind in {"pptx", "pptx-4x5"}:
        presentation = Presentation(path)
        assert len(presentation.slides) == 1
        if kind == "pptx-4x5":
            assert abs(
                presentation.slide_width / presentation.slide_height - 4 / 5
            ) < 0.001
        else:
            assert presentation.slide_width == presentation.slide_height
        slide = presentation.slides[0]
        text = "\n".join(
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
        )
        assert "Menos é mais." in text
        assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
        return

    if kind == "docx":
        document = Document(path)
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        assert "Relatório do mês" in text
        assert "Um documento simples produzido dentro dos trilhos da marca." in text
        assert {"Brand Heading", "Brand Body"}.issubset(
            {style.name for style in document.styles}
        )
        assert len(document.inline_shapes) >= 1
        return

    raise ValueError(f"Formato desconhecido: {kind}")


if __name__ == "__main__":
    main()
